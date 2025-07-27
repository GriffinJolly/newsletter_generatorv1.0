from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from datetime import datetime
import logging
import os
import json
import re
import random

# pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE, MSO_THEME_COLOR

# Ensure all pptx utilities are available
try:
    from pptx.util import Inches, Pt, Cm
except ImportError as e:
    logging.error(f"Failed to import pptx utilities: {e}")
    raise

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('newsletter_generator.log', mode='w', encoding='utf-8')
    ]
)
logger = logging.getLogger(__name__)

# Reduce verbosity for some noisy loggers
logging.getLogger('PIL').setLevel(logging.WARNING)
logging.getLogger('matplotlib').setLevel(logging.WARNING)

class NewsletterPPTGenerator:
    """Generate newsletter-style PowerPoint presentations from article insights"""
    
    def __init__(self, config: dict):
        """Initialize the PPT generator with configuration"""
        self.config = config
        self.template_path = Path(config.get('template_path', ''))
        self.output_dir = Path(config.get('output_dir', 'output'))
        self.logo_path = Path(config.get('logo_path', 'assets/logo.png'))
        
        # Create output directory if it doesn't exist
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Professional color palette
        self.colors = {
            'primary': (13, 59, 102),      # Dark blue
            'secondary': (27, 85, 226),    # Vibrant blue
            'accent': (1, 180, 228),       # Cyan accent
            'success': (57, 181, 74),      # Fresh green
            'warning': (255, 183, 3),      # Amber
            'background': (250, 250, 252), # Off-white
            'white': (255, 255, 255),      # Pure white
            'dark': (20, 24, 36),          # Near black
            'light_text': (108, 117, 125), # Gray text
            'border': (233, 236, 239),     # Light border
            'gradient_start': (245, 247, 250),  # Light gray gradient start
            'gradient_end': (255, 255, 255)     # White gradient end
        }
        
        # Typography settings with fallbacks for cross-platform compatibility
        self.fonts = {
            'heading': 'Montserrat, Arial, sans-serif',
            'subheading': 'Open Sans, Arial, sans-serif',
            'body': 'Segoe UI, Arial, sans-serif',
            'accent': 'Georgia, Times New Roman, serif',
            'mono': 'Courier New, monospace'
        }
        
        # Slide layouts
        self.slide_layouts = {
            'title': 0,
            'section_header': 1,
            'content': 2,
            'two_content': 3,
            'comparison': 4,
            'title_only': 5,
            'blank': 6,
            'content_with_caption': 7,
            'picture_with_caption': 8
        }
    
    def _create_new_presentation(self):
        """Create a new presentation with default settings"""
        from pptx import Presentation
        from pptx.util import Inches
        
        # Create a new presentation with a blank layout
        prs = Presentation()
        
        # Set slide size to 16:9 (widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Return the presentation
        return prs
        
    def _add_blank_slide(self, prs):
        """Add a blank slide to the presentation"""
        # Use the blank layout (index 6)
        blank_slide_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_slide_layout)
        
    def _add_newsletter_header(self, slide, title: str, issue_info: str = "") -> None:
        """Add a clean, structured newsletter header
        
        Args:
            slide: The slide object to add the header to
            title: The title text for the header
            issue_info: Optional issue information (e.g., date, issue number)
        """
        try:
            # Get slide width, defaulting to standard 16:9 width if not available
            slide_width = getattr(slide, 'width', Inches(13.33).pt)
            
            # Add logo if available
            logo = None
            title_left = Inches(0.5)  # Default left position without logo
            
            if hasattr(self, 'logo_path') and self.logo_path and self.logo_path.exists():
                try:
                    logo = slide.shapes.add_picture(
                        str(self.logo_path),
                        Inches(0.5), Inches(0.25),
                        height=Inches(1.0)
                    )
                    # Position title to the right of the logo
                    title_left = logo.left + logo.width + Inches(0.5)
                except Exception as e:
                    logger.warning(f"Could not add logo: {e}")
            
            # Calculate available width for title
            available_width = slide_width - title_left - Inches(0.5)
            
            # Add title
            title_box = slide.shapes.add_textbox(
                title_left, Inches(0.4),
                available_width, Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.name = self.fonts['heading'].split(',')[0].strip()
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.colors['primary'])  # Changed to primary color for better contrast
            p.space_after = 0
            
            # Add issue info if provided
            if issue_info:
                info_box = slide.shapes.add_textbox(
                    title_left, Inches(0.9),
                    available_width, Inches(0.4)
                )
                tf = info_box.text_frame
                p = tf.paragraphs[0]
                p.text = issue_info.upper()
                p.font.name = self.fonts['subheading'].split(',')[0].strip()
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(*self.colors['light_text'])
                p.space_after = 0
            
            # Add a subtle separator line
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, Inches(1.5) - Pt(2),
                slide_width, Pt(2)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*self.colors['accent'])
            line.line.fill.background()
            
        except Exception as e:
            logger.error(f"Error in _add_newsletter_header: {str(e)}")
            raise
        
    def _create_gradient_background(self, slide) -> None:
        """Create a subtle gradient background"""
        # Add a rectangle shape for background
        left = 0
        top = 0
        width = slide.parent.slide_width
        height = slide.parent.slide_height
        
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        
        # Set fill to gradient
        fill = bg_shape.fill
        fill.gradient()
        fill.gradient_angle = 45
        
        # Gradient stops
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[0].color.rgb = RGBColor(*self.colors['gradient_start'])
        fill.gradient_stops[1].position = 1.0
        fill.gradient_stops[1].color.rgb = RGBColor(*self.colors['gradient_end'])
        
        # Remove border
        bg_shape.line.fill.background()
        
        # Send to back
        bg_shape.z_order = 0
    
    def _add_decorative_header_bar(self, slide) -> None:
        """Add a modern header bar with subtle shadow"""
        # Main header bar
        left = 0
        top = 0
        width = slide.parent.slide_width
        height = Inches(0.25)
        
        # Add subtle shadow effect (semi-transparent black)
        shadow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            left + Pt(2), top + Pt(2), 
            width, height
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
        shadow.fill.alpha = 0.1
        shadow.line.fill.background()
        
        # Main header bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        
        # Gradient fill
        fill = header_bar.fill
        fill.gradient()
        fill.gradient_angle = 0
        
        # Gradient from primary to secondary color
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[0].color.rgb = RGBColor(*self.colors['primary'])
        fill.gradient_stops[1].position = 1.0
        fill.gradient_stops[1].color.rgb = RGBColor(*self.colors['secondary'])
        
        # Remove border
        header_bar.line.fill.background()
        
        # Add a thin accent line at the bottom
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top + height - Pt(2),
            width, Pt(4)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = RGBColor(*self.colors['accent'])
        accent_line.line.fill.background()
    
        # Create a header box
        header_box = self._create_content_box(
            slide, 
            (0.25, 0.25, 13.5, 1.5),
            has_border=False
        )
        
        # Add logo if available
        if self.logo_path.exists():
            try:
                logo_left = Inches(0.5)
                logo_top = Inches(0.4)
                logo = slide.shapes.add_picture(
                    str(self.logo_path), 
                    logo_left, logo_top, 
                    height=Inches(0.7)
                )
                title_left = Inches(1.8)
            except Exception as e:
                logger.warning(f"Could not add logo: {e}")
                title_left = Inches(0.5)
        else:
            title_left = Inches(0.5)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(title_left), Inches(0.4),
            Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        
        # Main title
        p = title_frame.paragraphs[0]
        p.text = title.upper()
        p.font.name = self.fonts['heading'].split(',')[0].strip()
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.colors['primary'])
        p.space_after = Pt(4)
        
        # Add issue info
        if issue_info:
            p = title_frame.add_paragraph()
            p.text = issue_info.upper()
            p.font.name = self.fonts['subheading'].split(',')[0].strip()
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(*self.colors['light_text'])
            p.space_before = 0
            p.space_after = 0
        
        # Add date on the right side
        if issue_info:
            date_box = slide.shapes.add_textbox(
                Inches(10), Inches(0.6),
                Inches(3), Inches(0.5)
            )
            date_frame = date_box.text_frame
            p = date_frame.paragraphs[0]
            p.text = datetime.now().strftime('%B %d, %Y')
            p.font.name = self.fonts['body'].split(',')[0].strip()
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(*self.colors['light_text'])
            p.alignment = PP_ALIGN.RIGHT
    
    def _create_content_box(self, slide, position: Tuple[float, float, float, float], 
                          title: str = "", has_border: bool = True) -> tuple:
        """Create a structured content box with optional title
        
        Args:
            slide: The slide to add the box to
            position: Tuple of (left, top, width, height) in inches
            title: Optional title for the content box
            has_border: Whether to draw a border around the box
            
        Returns:
            The created shape object
        """
        # Import pptx utilities locally to ensure they're in scope
        from pptx.util import Inches, Pt, Cm
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        
        left, top, width, height = position
        
        try:
            # Convert to points for precision if not already in points
            left_pt = left if isinstance(left, int) else Inches(left)
            top_pt = top if isinstance(top, int) else Inches(top)
            width_pt = width if isinstance(width, int) else Inches(width)
            height_pt = height if isinstance(height, int) else Inches(height)
        except Exception as e:
            logger.error(f"Error converting dimensions: {e}")
            logger.error(f"Left: {left} (type: {type(left)}), Top: {top} (type: {type(top)}), "
                       f"Width: {width} (type: {type(width)}), Height: {height} (type: {type(height)})")
            raise
        
        # Create main box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_pt, top_pt,
            width_pt, height_pt
        )
        
        # Style the box
        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.colors['white'])
        
        line = box.line
        if has_border:
            line.color.rgb = RGBColor(*self.colors['border'])
            line.width = Pt(1)
        else:
            line.fill.background()
        
        # Add title if provided
        title_box = None
        if title:
            # Title background
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_pt, top_pt,
                width_pt, Pt(40)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(*self.colors['primary'])
            title_bg.line.fill.background()
            
            # Title text
            title_box = slide.shapes.add_textbox(
                left_pt + Pt(10), top_pt + Pt(5),
                width_pt - Pt(20), Pt(30)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.name = self.fonts['heading'].split(',')[0].strip()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.colors['white'])
            p.space_after = 0
            
            # Adjust content area
            top_pt += Pt(45)
            height_pt -= Pt(45)
        
        # Return content area coordinates
        from pptx.util import Inches, Pt
        from textwrap import wrap

        # Convert position tuple to inches
        left, top, width, height = position
        left_pt = Inches(left)
        top_pt = Inches(top)
        width_pt = Inches(width)
        height_pt = Inches(height)

        # Add a card with shadow
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pt, top_pt, width_pt, height_pt
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        card.line.color.rgb = RGBColor(230, 230, 230)  # Light gray border

        # Add shadow effect (using default shadow color)
        shadow = card.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = Pt(4)
        shadow.offset_x = Pt(2)
        shadow.offset_y = Pt(2)

        # Content will be added by the calling method
        # (e.g., _add_article_card)
        return box

        # Add content with padding
        padding = Pt(12)
        content_left = left_pt + padding
        content_top = top_pt + padding + Pt(8)  # Extra space for accent
        content_width = width_pt - (2 * padding)

        # Add title with better typography
        title = str(article.get('title', 'Untitled Article'))
        title_box = slide.shapes.add_textbox(
            content_left, content_top,
            content_width, Pt(40)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(16)  # Slightly larger for better readability
        p.font.bold = True
        p.font.color.rgb = RGBColor(20, 30, 40)
        p.space_after = Pt(6)

        # Add source and date with better styling
        source = article.get('source', 'Source')
        date = article.get('date', '')
        source_text = f"{source.upper()} • {date}"
        source_box = slide.shapes.add_textbox(
            content_left, content_top + Pt(28),
            content_width, Pt(15)
        )
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = source_text
        p.font.name = 'Arial'
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(1, 180, 228)  # Accent color
        p.space_after = Pt(12)

        # Ensure summary is not truncated to 2 sentences or with ellipsis
        tf = title_box.text_frame
        # Take up to 15 sentences or 1200 chars, ending at a sentence boundary
        import re
        sentences = re.split(r'(?<=[.!?])\s+', detailed_summary)
        summary_text = ' '.join(sentences[:15])
        if len(summary_text) > 1200:
            # Trim to 1200 chars, but end at last period
            cut = summary_text[:1200]
            last_period = cut.rfind('.')
            if last_period > 0:
                summary_text = cut[:last_period+1]
            else:
                summary_text = cut
        # Split into paragraphs for better readability
        paragraphs = summary_text.split('\n\n')
        for i, para in enumerate(paragraphs):
            if i > 0:
                p = tf.add_paragraph()
                p.text = ''  # Add space between paragraphs
                p.space_after = Pt(4)
            p = tf.add_paragraph()
            p.text = para.strip()
            p.font.name = 'Arial'
            p.font.size = Pt(10)
            p.space_after = Pt(4)
            p.space_before = Pt(0)

        # Add key points section if available
        key_points = article.get('key_points', [])
        if key_points:
            # Add key points header
            kp_header = slide.shapes.add_textbox(
                content_left, content_top + Pt(50),
                content_width, Pt(15)
            )
            tf = kp_header.text_frame
            p = tf.paragraphs[0]
            p.text = "KEY POINTS"
            p.font.name = 'Arial'
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.space_after = Pt(2)
            
            # Add key points content
            kp_box = slide.shapes.add_textbox(
                content_left, content_top + Pt(65),
                content_width, Pt(100)
            )
            tf = kp_box.text_frame
            
            for i, point in enumerate(key_points[:5]):  # Limit to 5 key points
                if i > 0:
                    p = tf.add_paragraph()
                    p.text = ''
                    p.space_after = Pt(2)
                
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.name = 'Arial'
                p.font.size = Pt(9)
                p.space_after = Pt(2)
                p.level = 1
        
        # Add entities section if available
        entities = article.get('entities', {})
        if entities and any(entities.values()):
            # Add entities header
            ent_header = slide.shapes.add_textbox(
                content_left, content_top + Pt(165),
                content_width, Pt(12)
            )
            tf = ent_header.text_frame
            p = tf.paragraphs[0]
            p.text = "MENTIONED"
            p.font.name = 'Arial'
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = RGBColor(100, 100, 100)
            
            # Add entities content
            ent_box = slide.shapes.add_textbox(
                content_left, content_top + Pt(175),
                content_width, Pt(50)
            )
            tf = ent_box.text_frame
            
            # Format entities by type
            entity_texts = []
            for ent_type in ['organizations', 'people', 'locations']:
                if entities.get(ent_type):
                    entities_list = [e for e in entities[ent_type] if e][:3]  # Limit to 3 per type
                    if entities_list:
                        entity_texts.append(f"{ent_type.title()}: {', '.join(entities_list)}")
            
            if entity_texts:
                p = tf.paragraphs[0]
                p.text = ' | '.join(entity_texts)
                p.font.name = 'Arial'
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(120, 120, 120)
        
        # Add category tags at bottom if available
        if 'categories' in article and article['categories']:
            categories = article['categories'][:2]  # Show max 2 categories
            tag_top = top_pt + height_pt - Pt(25)
            
            # Add a subtle separator line above tags
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                content_left, tag_top - Pt(10),
                content_width, Pt(0.5)  # Thin rectangle as a line
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(230, 230, 230)
            line.line.fill.background()  # No border
            
            for i, category in enumerate(categories):
                if i > 1:  # Only show first 2 categories
                    break
                    
                # Create tag background
                tag_width = min(Inches(1.5), Inches(len(category) * 0.15))
                tag_left = content_left + (Inches(1.8) * i)
                
                tag_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    tag_left, tag_top - Pt(5),
                    tag_width, Pt(18)
                )
                tag_bg.fill.solid()
                tag_bg.fill.fore_color.rgb = RGBColor(240, 245, 250)  # Light blue-gray
                tag_bg.line.color.rgb = RGBColor(200, 220, 240)
                tag_bg.line.width = Pt(0.5)
                
                # Add tag text
                tag_text = slide.shapes.add_textbox(
                    tag_left + Pt(4), tag_top - Pt(5),
                    tag_width - Pt(8), Pt(18)
                )
                tf = tag_text.text_frame
                p = tf.paragraphs[0]
                p.text = category.upper()
                p.font.name = 'Arial'
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = RGBColor(70, 100, 150)
                p.alignment = PP_ALIGN.CENTER
    
    def _enhance_summary(self, summary: str, article: Dict[str, Any]) -> str:
        """Generate a comprehensive summary of at least 15 sentences.
        
        Args:
            summary: The original summary text
            article: Dictionary containing article data
            
        Returns:
            str: Enhanced summary text
        """
        try:
            content = str(article.get('content', ''))
            if not content:
                logger.warning(f"Article has no content: {article.get('title', 'Untitled')}")
                if 'url' in article:
                    logger.warning(f"Article URL: {article['url']}")
                return summary if summary and len(summary) > 20 else "No sufficient content available for this article."
                
            # Split content into sentences using multiple delimiters
            logger.debug("Splitting content into sentences...")
            # Use a fixed-width lookbehind regex for sentence splitting
            sentences = re.split(r'(?<=[.!?])\s+', content)
            logger.debug(f"Found {len(sentences)} initial sentences")
            
            # Log first few sentences for debugging
            for i, s in enumerate(sentences[:5]):
                logger.debug(f"Sentence {i+1} (len: {len(s)}): {s[:100]}{'...' if len(s) > 100 else ''}")
            
            # Clean and filter sentences
            cleaned_sentences = []
            short_sentences = 0
            long_sentences = 0
            
            for sent in sentences:
                sent = sent.strip()
                words = sent.split()
                word_count = len(words)
                
                # Log sentence statistics
                if word_count < 5:
                    short_sentences += 1
                    continue
                elif word_count > 50:
                    long_sentences += 1
                    # Try to split very long sentences
                    parts = re.split(r'[,;:]\s+', sent)
                    if len(parts) > 1 and all(5 <= len(p.split()) <= 30 for p in parts):
                        cleaned_sentences.extend(parts)
                        continue
                
                # Include sentences between 5-50 words
                if 5 <= word_count <= 50:
                    cleaned_sentences.append(sent)
                
                if len(cleaned_sentences) >= 25:  # Target 25 sentences
                    break
            
            logger.debug(f"Filtered sentences: {len(cleaned_sentences)} (short: {short_sentences}, long: {long_sentences})")
            
            # If we have enough sentences, use them
            if len(cleaned_sentences) >= 15:
                result = ' '.join(cleaned_sentences[:15])
                logger.debug(f"Generated summary length: {len(result)} characters, {len(result.split())} words")
                return result
            elif len(cleaned_sentences) > 0:
                # If not enough, use as many as possible, then fallback below
                partial = ' '.join(cleaned_sentences)
            else:
                partial = ''
                
            # Fallback: Try to get more content using different strategies
            logger.debug("Trying fallback content extraction...")
            
            # Strategy 1: Use the partial summary and fill with more from content
            if partial:
                fallback_text = content[len(partial):][:1000]
                fallback_sentences = re.split(r'(?<=[.!?])\s+', fallback_text)
                needed = 15 - len(partial.split('. '))
                combined = partial + ' ' + ' '.join(fallback_sentences[:needed])
                result = combined.strip()
                logger.debug(f"Fallback filled summary length: {len(result)} characters")
                return result
            
            # Strategy 2: Try to extract paragraphs
            paragraphs = [p for p in content.split('\n\n') if len(p.split()) > 10]
            if paragraphs:
                result = ' '.join(paragraphs[:5])  # Join first 5 paragraphs
                logger.debug(f"Fallback 2 (paragraphs) summary length: {len(result)} characters")
                return result
                
            # Last resort: return first 2000 characters with proper truncation
            fallback_text = content[:2000]
            last_period = fallback_text.rfind('.')
            if last_period > 500:
                result = fallback_text[:last_period + 1]
            else:
                result = fallback_text
            logger.warning(f"Using last-resort content extraction. Length: {len(result)} characters")
            return result
            
        except Exception as e:
            logger.error(f"Error enhancing summary: {str(e)}")
            return summary if summary else "No summary available."

    def _add_title_slide(self, prs, title: str, subtitle: str = "") -> None:
        """Add a title slide to the presentation"""
        # Use the title layout (index 0)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title
        title_box = slide.shapes.title
        title_box.text = title
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.font.name = self.fonts['heading'].split(',')[0].strip()
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.colors['primary'])
        
        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.placeholders[1]
            subtitle_box.text = subtitle
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
        p.font.name = self.fonts['heading'].split(',')[0].strip()
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.colors['primary'])
        
    def _add_article_card(self, slide, article: Dict[str, Any], position: Tuple[float, float, float, float]) -> bool:
        """Add an article card to the slide
        
        Args:
            slide: The slide to add the card to
            article: Dictionary containing article data
            position: Tuple of (left, top, width, height) in inches
        """
        try:
            left, top, width, height = position
            
            # Create a content box for the article card
            self._create_content_box(
                slide,
                (left, top, width, height),
                has_border=True
            )
            
            # Add article title
            title_left = left + 0.2
            title_top = top + 0.2
            title_width = width - 0.4
            
            title_box = slide.shapes.add_textbox(
                Inches(title_left), Inches(title_top),
                Inches(title_width), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = article.get('title', 'No Title')
            p.font.name = self.fonts['heading'].split(',')[0].strip()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.colors['primary'])
            p.space_after = Pt(6)
            
            # Add article source and date
            source_text = f"{article.get('source', 'Unknown')} • {article.get('date', '')}"
            source_box = slide.shapes.add_textbox(
                Inches(title_left), Inches(top + 1.0),
                Inches(title_width), Inches(0.3)
            )
            tf = source_box.text_frame
            p = tf.paragraphs[0]
            p.text = source_text.upper()
            p.font.name = self.fonts['subheading'].split(',')[0].strip()
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(*self.colors['light_text'])
            p.space_after = Pt(6)
            
            # Add article summary/content
            summary = article.get('summary', '')
            content = article.get('content', '')
            if content and len(content) > len(summary):
                display_content = content
            else:
                display_content = summary or 'No content available.'
            logger.info(f"[PPT] Article: '{article.get('title','')}' | Summary len: {len(summary)} | Content len: {len(content)} | Used: {len(display_content)}")
            # Filtering: skip if content is too short or irrelevant
            min_sentences = 3
            max_sentences = 15
            irrelevant_phrases = [
                "no content available", "click here", "read more", "subscribe", "copyright",
                "submitting email", "terms and conditions"
            ]
            
            # Remove sentences containing irrelevant phrases
            import re
            sentences = re.split(r'(?<=[.!?])\s+', display_content)
            filtered_sentences = [s for s in sentences if not any(phrase in s.lower() for phrase in irrelevant_phrases)]
            
            # If not enough sentences, try fallback (original summary/content)
            if len(filtered_sentences) < min_sentences:
                fallback = article.get('summary', '') or article.get('content', '')
                fallback_sentences = re.split(r'(?<=[.!?])\s+', fallback)
                fallback_filtered = [s for s in fallback_sentences if not any(phrase in s.lower() for phrase in irrelevant_phrases)]
                if len(fallback_filtered) >= min_sentences:
                    filtered_sentences = fallback_filtered
                    logger.info(f"[PPT] Used fallback content for article '{article.get('title','')}'")
                else:
                    logger.info(f"[PPT] Skipping article '{article.get('title','')}' due to insufficient/irrelevant content (even after fallback).")
                    return False  # Skip this slide, return False
            
            # Truncate to max sentences
            truncated_content = ' '.join(filtered_sentences[:max_sentences]).strip()
            if len(filtered_sentences) > max_sentences:
                truncated_content += ' ...'
            display_content = truncated_content
            content_box = slide.shapes.add_textbox(
                Inches(title_left), Inches(top + 1.4),
                Inches(title_width - 0.2), Inches(height - 1.6)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # Add content in paragraphs
            paragraphs = content.split('\n\n')
            for i, para in enumerate(paragraphs):
                if i > 0:
                    p = tf.add_paragraph()
                    p.space_before = Pt(6)
                else:
                    p = tf.paragraphs[0]
                p.text = para
                p.font.name = self.fonts['body'].split(',')[0].strip()
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(*self.colors['dark'])
                p.space_after = Pt(3)
                p.alignment = PP_ALIGN.LEFT
                
            # Add category tag if available
            categories = article.get('categories', [])
            if categories:
                category = categories[0]  # Use first category as tag
                tag_left = left + width - 1.5
                tag_top = top + height - 0.5
                
                tag = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(tag_left), Inches(tag_top),
                    Inches(1.3), Inches(0.3)
                )
                tag.fill.solid()
                tag.fill.fore_color.rgb = RGBColor(*self.colors['accent'])
                tag.line.fill.background()
                
                # Add category text
                tag_text = slide.shapes.add_textbox(
                    Inches(tag_left + 0.1), Inches(tag_top + 0.05),
                    Inches(1.1), Inches(0.2)
                )
                tf = tag_text.text_frame
                p = tf.paragraphs[0]
                p.text = category.upper()
                p.font.name = self.fonts['subheading'].split(',')[0].strip()
                p.font.size = Pt(7)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*self.colors['white'])
                p.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            logger.error(f"Error adding article card: {str(e)}")
            return False  # On error, return False
        return True  # If no error or skip, slide was added successfully
    
    def _add_summary_slide(self, prs, insights: List[Dict[str, Any]]) -> None:
        """Add a summary slide to the presentation"""
        # Use the content layout (index 2)
        content_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(content_layout)
        
        # Add title
        title_box = slide.shapes.title
        title_box.text = "Summary"
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.font.name = self.fonts['heading'].split(',')[0].strip()
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.colors['primary'])
        
        # Add summary content
        summary_box = slide.shapes.placeholders[1]
        # Combine summaries from all insights, then enhance
        combined_summary = '\n'.join([
            insight.get('summary', '') or insight.get('content', '')
            for insight in insights if insight.get('summary') or insight.get('content')
        ])
        # Use _enhance_summary to generate a detailed summary
        detailed_summary = self._enhance_summary(combined_summary, {'content': combined_summary})
        summary_box.text = detailed_summary
        tf = summary_box.text_frame
        # Split into paragraphs for better readability
        for i, para in enumerate(detailed_summary.split('\n')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = para.strip()
            p.font.name = self.fonts['body'].split(',')[0].strip()
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(*self.colors['dark'])
        
        # Add a chart to visualize insights
        chart_left = Inches(1)
        chart_top = Inches(3)
        chart_width = Inches(8)
        chart_height = Inches(4)
        
        # Chart title
        title_box = slide.shapes.add_textbox(
            chart_left, Inches(2.5), chart_width, Inches(0.4)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Insights by Category"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.colors['primary'])
        p.alignment = PP_ALIGN.CENTER
        
        # Create a simple bar chart representation
        chart_data = []
        for insight in insights:
            for category in insight.get('categories', ['Other']):
                chart_data.append((category, 1))
        
        # Count insights by category
        category_counts = {}
        for category, _ in chart_data:
            category_counts[category] = category_counts.get(category, 0) + 1
        
        # Create bars
        max_count = max(category_counts.values()) if category_counts else 1
        bar_height = 0.4
        bar_spacing = 0.6
        colors = [self.colors['primary'], self.colors['secondary'], self.colors['accent'], 
                 self.colors['success'], self.colors['warning']]
        
        y_pos = chart_top
        for i, (category, _) in enumerate(category_counts.items()):
            # Bar
            bar_width = (category_counts[category] / max_count) * (chart_width.inches - 2)
            
            # Create bar shape
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(chart_left.inches + 2), 
                y_pos,
                Inches(bar_width), 
                Inches(bar_height)
            )
            
            # Style the bar
            color = colors[i % len(colors)]
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*color)
            bar.line.fill.background()
            
            # Add category label
            label = slide.shapes.add_textbox(
                chart_left, 
                y_pos,
                Inches(1.8), 
                Inches(bar_height)
            )
            
            # Configure label text
            tf = label.text_frame
            p = tf.paragraphs[0]
            p.text = f"{category} ({category_counts[category]})"
            p.font.name = 'Arial'
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(*self.colors['dark'])
            p.alignment = PP_ALIGN.LEFT
            
            # Position for next bar
            y_pos += Inches(bar_height + 0.2)
            
            # Stop if we're running out of vertical space
            if y_pos > Inches(6.5):
                break
                
        # Add a simple legend
        legend_left = Inches(8)
        legend_top = Inches(2)
        legend_size = Inches(0.3)
        
        for i, (category, _) in enumerate(category_counts.items()):
            if i >= 5:  # Limit legend items
                break
                
            # Color swatch
            swatch = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                legend_left,
                legend_top + (i * Inches(0.5)),
                legend_size,
                legend_size
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = RGBColor(*colors[i % len(colors)])
            swatch.line.fill.background()
            
            # Category name
            cat_label = slide.shapes.add_textbox(
                legend_left + legend_size + Inches(0.2),
                legend_top + (i * Inches(0.5)),
                Inches(2),
                legend_size
            )
            tf = cat_label.text_frame
            p = tf.paragraphs[0]
            p.text = category
            p.font.name = 'Arial'
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(*self.colors['dark'])
            
        # Add a call to action at the bottom
        cta_box = slide.shapes.add_textbox(
            Inches(1), Inches(6),
            slide.width - Inches(2), Inches(1)
        )
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Read more about these articles and stay up-to-date with the latest news!"
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    def generate_presentation(self, insights: List[Dict[str, Any]], output_path: Optional[str] = None) -> str:
        """
        Generate a PowerPoint presentation from the given insights.

        Args:
            insights: List of insight dictionaries
            output_path: Optional path to save the presentation

        Returns:
            str: Path to the generated presentation file
        """
        try:
            # Create a new presentation
            self.prs = self._create_new_presentation()
            
            # Group insights by sector
            insights_by_sector = {}
            for insight in insights:
                sector = insight.get('sector', 'General')
                if sector not in insights_by_sector:
                    insights_by_sector[sector] = []
                insights_by_sector[sector].append(insight)
            
            # Add title slide
            self._add_title_slide(self.prs, "Weekly Newsletter", datetime.now().strftime("%B %d, %Y"))
            
            # Prepare fetchers for replacement logic
            from scrapers.news_fetcher import NewsAPIFetcher, GNewsFetcher
            fetcher_classes = {'NewsAPIFetcher': NewsAPIFetcher, 'GNewsFetcher': GNewsFetcher}
            
            # Add section and article slides
            for sector, sector_insights in insights_by_sector.items():
                # Add section header
                section_slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.slide_layouts['section_header']])
                self._add_newsletter_header(section_slide, sector)
                
                # Positioning for article cards
                n_articles = len(sector_insights)
                max_per_slide = 2
                slide = None
                pos_idx = 0
                positions = [
                    (0.5, 1.5, 6, 4.5),
                    (7, 1.5, 6, 4.5)
                ]
                
                used_urls = set()
                for idx, article in enumerate(sector_insights):
                    if pos_idx == 0:
                        # Start new slide
                        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.slide_layouts['content']])
                        self._add_decorative_header_bar(slide)
                    
                    url = article.get('url')
                    if url:
                        used_urls.add(url)
                    added = self._add_article_card(slide, article, positions[pos_idx])
                    
                    # If not added, try replacement logic
                    if not added:
                        # Determine which fetcher to use for replacement
                        orig_source = article.get('source', '').lower()
                        # Prefer alternate fetcher (not WebScraper, not SECFetcher)
                        fetcher_order = ['NewsAPIFetcher', 'GNewsFetcher']
                        if 'gnews' in orig_source:
                            fetcher_order = ['NewsAPIFetcher']
                        elif 'newsapi' in orig_source:
                            fetcher_order = ['GNewsFetcher']
                        else:
                            fetcher_order = ['NewsAPIFetcher', 'GNewsFetcher']
                        replacement_added = False
                        for fetcher_name in fetcher_order:
                            try:
                                fetcher = fetcher_classes[fetcher_name](self.config)
                                # Use sector as query; fallback to article title if sector is generic
                                query = sector if sector.lower() != 'general' else article.get('title', '')
                                results = fetcher.fetch(query, max_results=5)
                                for cand in results:
                                    cand_url = cand.get('url')
                                    if cand_url and cand_url in used_urls:
                                        continue
                                    # Mark as used even if not added, to avoid infinite loops
                                    if cand_url:
                                        used_urls.add(cand_url)
                                    cand_added = self._add_article_card(slide, cand, positions[pos_idx])
                                    if cand_added:
                                        logger.info(f"[PPT] Replacement article used from {fetcher_name} for sector '{sector}'")
                                        replacement_added = True
                                        break
                                if replacement_added:
                                    break
                            except Exception as e:
                                logger.error(f"[PPT] Error fetching replacement from {fetcher_name}: {e}")
                        if not replacement_added:
                            logger.info(f"[PPT] No suitable replacement found for skipped article in sector '{sector}'")
                    pos_idx += 1
                    if pos_idx == max_per_slide or idx == n_articles - 1:
                        pos_idx = 0
            
            # Add summary slide at the end
            self._add_summary_slide(self.prs, insights)
            
            # Save the presentation
            if not output_path:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = self.output_dir / f"newsletter_{timestamp}.pptx"
            else:
                output_path = Path(output_path)
            
            # Ensure output directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            self.prs.save(str(output_path))
            logger.info(f"Newsletter presentation saved to {output_path}")
            
            return str(output_path)
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise


def generate_sample_insights() -> List[Dict[str, Any]]:
    """Generate sample insights for testing the newsletter format"""
    return [
        {
            "title": "AI Revolution Transforms Healthcare Diagnostics",
            "content": "Revolutionary AI technology is enabling faster and more accurate medical diagnoses, with new algorithms showing 95% accuracy in detecting early-stage diseases. The technology combines machine learning with advanced imaging to provide real-time analysis that could save millions of lives annually.",
            "source": "MedTech Today",
            "date": "2024-06-15",
            "categories": ["Healthcare", "Artificial Intelligence", "Innovation"],
            "sector": "Healthcare Technology",
            "key_points": [
                "95% accuracy in early disease detection",
                "Real-time analysis capabilities",
                "Potential to save millions of lives"
            ]
        },
        {
            "title": "Green Energy Breakthrough: Solar Efficiency Reaches New Heights",
            "content": "Scientists have achieved a major breakthrough in solar panel technology, reaching 47% efficiency in laboratory conditions. This advancement could revolutionize renewable energy adoption and significantly reduce costs for consumers worldwide.",
            "source": "Energy Innovation Weekly",
            "date": "2024-06-12",
            "categories": ["Renewable Energy", "Technology", "Sustainability"],
            "sector": "Clean Energy",
            "key_points": [
                "47% solar panel efficiency achieved",
                "Could revolutionize renewable energy",
                "Significant cost reduction potential"
            ]
        },
        {
            "title": "Quantum Computing Makes Commercial Debut",
            "content": "The first commercial quantum computer has been deployed for financial modeling, promising to solve complex optimization problems in seconds rather than hours. This marks a significant milestone in quantum technology commercialization.",
            "source": "Quantum Business Review",
            "date": "2024-06-10",
            "categories": ["Quantum Computing", "Finance", "Technology"],
            "sector": "Quantum Technology",
            "key_points": [
                "First commercial quantum deployment",
                "Solves complex problems in seconds",
                "Major milestone for commercialization"
            ]
        },
        {
            "title": "Space Tourism Industry Reaches New Milestone",
            "content": "Commercial space flights have successfully transported over 1,000 passengers to the edge of space, marking a significant achievement for the space tourism industry. The milestone demonstrates growing consumer confidence and technological maturity.",
            "source": "Space Commerce Daily",
            "date": "2024-06-08",
            "categories": ["Space Tourism", "Transportation", "Innovation"],
            "sector": "Aerospace",
            "key_points": [
                "Over 1,000 passengers transported",
                "Growing consumer confidence",
                "Technological maturity demonstrated"
            ]
        },
        {
            "title": "Autonomous Vehicles Pass Major Safety Milestone",
            "content": "Self-driving cars have achieved a 99.9% safety record in controlled testing environments, surpassing human driver performance in multiple categories. The breakthrough brings fully autonomous transportation closer to reality.",
            "source": "AutoTech Insider",
            "date": "2024-06-05",
            "categories": ["Autonomous Vehicles", "Safety", "Transportation"],
            "sector": "Automotive Technology",
            "key_points": [
                "99.9% safety record achieved",
                "Surpasses human driver performance",
                "Brings full autonomy closer to reality"
            ]
        }
    ]


# Example usage
if __name__ == "__main__":
    # Configuration
    config = {
        'output_dir': 'newsletter_output',
        'logo_path': 'assets/company_logo.png'
    }
    
    # Create generator
    generator = NewsletterPPTGenerator(config)
    
    # Generate sample presentation
    sample_insights = generate_sample_insights()
    output_file = generator.generate_presentation(sample_insights)
    
    print(f"Newsletter presentation generated: {output_file}")